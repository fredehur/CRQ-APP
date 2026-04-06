# tools/export_pptx.py
"""
Builds the board PPTX report using python-pptx.
Bootstraps tools/templates/base.pptx on first run if absent.

Usage:
    uv run python tools/export_pptx.py [output.pptx]
    Defaults to output/board_report.pptx
"""
from __future__ import annotations

import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

sys.path.insert(0, os.path.dirname(__file__))
sys.path.insert(0, str(Path(__file__).parent.parent))
from report_builder import build, ReportData, RegionEntry, RegionStatus
from tools.config import BOARD_PPTX_PATH

TEMPLATES_DIR = Path(__file__).parent / "templates"
BASE_PPTX     = TEMPLATES_DIR / "base.pptx"
DEFAULT_OUT   = str(BOARD_PPTX_PATH)

# Brand colours
BRAND   = RGBColor(0x10, 0x42, 0x77)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
SLATE   = RGBColor(0x64, 0x74, 0x8B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1E, 0x29, 0x3B)

# Slide dimensions (standard 10×7.5in widescreen 16:9)
W = Inches(10)
H = Inches(7.5)


# ── Base PPTX bootstrap ────────────────────────────────────────────────────────

def _ensure_base_pptx() -> None:
    """Generate tools/templates/base.pptx if it does not exist."""
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    if BASE_PPTX.exists():
        return

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    # Remove all default slide layouts except blank (index 6)
    # python-pptx requires at least one layout — keep blank
    prs.save(str(BASE_PPTX))


# ── Slide helpers ──────────────────────────────────────────────────────────────

def _add_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(blank_layout)


def _fill_shape(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, width, height
    )
    shape.line.fill.background()
    _fill_shape(shape, color)
    return shape


def _add_text(slide, text: str, left, top, width, height,
              font_size: int = 14, bold: bool = False,
              color: RGBColor = DARK, align=PP_ALIGN.LEFT,
              wrap: bool = True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _status_label(status: RegionStatus, dominant_pillar: str | None = None) -> str:
    if status == RegionStatus.ESCALATED:
        if dominant_pillar == "Geopolitical": return "ESCALATED — GEO-LED"
        if dominant_pillar == "Cyber":        return "ESCALATED — CYBER-LED"
        return "ESCALATED"
    if status == RegionStatus.MONITOR: return "MONITOR"
    if status == RegionStatus.CLEAR:   return "CLEAR"
    return "UNKNOWN"


def _vel_label(velocity: str | None) -> str:
    mapping = {"accelerating": "↑ Accelerating", "improving": "↓ Improving",
               "stable": "→ Stable", "unknown": "—"}
    return mapping.get(velocity or "unknown", velocity or "—")


# ── Slide builders ─────────────────────────────────────────────────────────────

def build_cover(prs: Presentation, data: ReportData) -> None:
    slide = _add_blank_slide(prs)

    # Brand header band (top 45%)
    band_h = Inches(3.4)
    _add_rect(slide, 0, 0, W, band_h, BRAND)
    _add_text(slide, "AeroGrid Wind Solutions",
              Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.5),
              font_size=13, color=RGBColor(0xBF,0xDB,0xFF))
    _add_text(slide, "Global Cyber Risk\nIntelligence Brief",
              Inches(0.6), Inches(1.1), Inches(8.8), Inches(2.0),
              font_size=28, bold=True, color=WHITE)

    # Status strip — ESCALATED / MONITOR / CLEAR counts
    strip_data = [
        (str(data.escalated_count), "ESCALATED", RED),
        (str(data.monitor_count),   "MONITOR",   AMBER),
        (str(data.clear_count),     "CLEAR",     GREEN),
    ]
    for i, (count, label, color) in enumerate(strip_data):
        x = Inches(0.6 + i * 1.5)
        _add_rect(slide, x, Inches(3.7), Inches(1.2), Inches(0.75), color)
        _add_text(slide, count, x, Inches(3.73), Inches(1.2), Inches(0.4),
                  font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(slide, label, x, Inches(4.1), Inches(1.2), Inches(0.3),
                  font_size=7, color=WHITE, align=PP_ALIGN.CENTER)

    # Meta
    meta = f"Pipeline ID: {data.run_id}   |   {data.timestamp}"
    _add_text(slide, meta, Inches(0.6), Inches(5.0), Inches(8), Inches(0.4),
              font_size=9, color=SLATE)

    # Confidential
    _add_text(slide, "CONFIDENTIAL",
              Inches(7.5), Inches(6.9), Inches(2.3), Inches(0.4),
              font_size=8, bold=True, color=RED, align=PP_ALIGN.RIGHT)


def build_exec_summary(prs: Presentation, data: ReportData) -> None:
    slide = _add_blank_slide(prs)

    # Header band
    _add_rect(slide, 0, 0, W, Inches(0.65), BRAND)
    _add_text(slide, "Executive Summary", Inches(0.3), Inches(0.1),
              Inches(7), Inches(0.45), font_size=16, bold=True, color=WHITE)

    # Status badges — no VaCR badge
    badge_data = [
        (str(data.escalated_count), "ESCALATED", RED),
        (str(data.monitor_count),   "MONITOR",   AMBER),
        (str(data.clear_count),     "CLEAR",     GREEN),
    ]
    for i, (count, label, color) in enumerate(badge_data):
        x = Inches(0.4 + i * 1.3)
        _add_rect(slide, x, Inches(0.85), Inches(1.1), Inches(0.75), color)
        _add_text(slide, count,  x, Inches(0.88), Inches(1.1), Inches(0.4),
                  font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(slide, label, x, Inches(1.25), Inches(1.1), Inches(0.3),
                  font_size=7, color=WHITE, align=PP_ALIGN.CENTER)

    # Exec summary text
    _add_text(slide, data.exec_summary,
              Inches(0.4), Inches(1.8), Inches(9.2), Inches(1.0),
              font_size=10, color=DARK, wrap=True)

    # Region table header — no VaCR, no Admiralty; add Confidence
    cols = ["Region", "Scenario", "Severity", "Velocity", "Confidence"]
    col_widths = [Inches(1.2), Inches(2.2), Inches(1.4), Inches(1.8), Inches(1.4)]
    y_header = Inches(3.0)
    x = Inches(0.4)
    for col, w in zip(cols, col_widths):
        _add_text(slide, col, x, y_header, w, Inches(0.3),
                  font_size=8, bold=True, color=SLATE)
        x += w

    escalated = [r for r in data.regions if r.status == RegionStatus.ESCALATED]
    for row_i, r in enumerate(escalated):
        y = Inches(3.35 + row_i * 0.45)
        row = [r.name, r.scenario_match or "—", r.severity or "—",
               _vel_label(r.velocity), r.confidence_label or "—"]
        x = Inches(0.4)
        for val, w in zip(row, col_widths):
            _add_text(slide, val, x, y, w, Inches(0.35), font_size=9, color=DARK)
            x += w


def build_region(prs: Presentation, region: RegionEntry) -> None:
    slide = _add_blank_slide(prs)

    # Header band
    _add_rect(slide, 0, 0, W, Inches(0.65), BRAND)
    _add_text(slide, region.name, Inches(0.3), Inches(0.1),
              Inches(6), Inches(0.45), font_size=16, bold=True, color=WHITE)
    _add_text(slide, region.scenario_match or "",
              Inches(0.3), Inches(0.38), Inches(6), Inches(0.25),
              font_size=9, color=RGBColor(0xBF,0xDB,0xFF))

    # Severity chip — colour by severity
    sev = (region.severity or "").upper()
    if sev == "CRITICAL":
        chip_color = RED
    elif sev == "HIGH":
        chip_color = RGBColor(0xEA, 0x58, 0x0C)
    elif sev == "MEDIUM":
        chip_color = AMBER
    else:
        chip_color = SLATE
    chip_label = region.severity or "—"
    _add_rect(slide, Inches(8.4), Inches(0.12), Inches(1.3), Inches(0.38), chip_color)
    _add_text(slide, chip_label, Inches(8.4), Inches(0.16),
              Inches(1.3), Inches(0.3), font_size=8, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)

    # Status label + meta row
    status_label = _status_label(region.status, region.dominant_pillar)
    meta_parts = [status_label]
    if region.threat_actor:
        meta_parts.append(f"Threat Actor: {region.threat_actor}")
    if region.signal_type_label:
        meta_parts.append(region.signal_type_label)
    if region.confidence_label:
        meta_parts.append(f"Confidence: {region.confidence_label}")
    _add_text(slide, "  ·  ".join(meta_parts),
              Inches(0.3), Inches(0.75), Inches(9), Inches(0.28),
              font_size=9, color=SLATE)

    # 5 intelligence rows (vertical list)
    rows = [
        ("INTEL FINDINGS",       region.intel_bullets[0]     if region.intel_bullets     else None, BRAND),
        ("ADVERSARY ACTIVITY",   region.adversary_bullets[0] if region.adversary_bullets else None, BRAND),
        ("IMPACT FOR AEROGRID",  region.impact_bullets[0]    if region.impact_bullets    else None, RED),
        ("WATCH FOR",            region.watch_bullets[0]     if region.watch_bullets     else None, AMBER),
        ("RECOMMENDED ACTION",   region.action_bullets[0]    if region.action_bullets    else None, SLATE),
    ]

    y = Inches(1.1)
    row_h = Inches(1.0)
    for label, text, color in rows:
        if not text:
            y += row_h
            continue
        # Left border stripe
        _add_rect(slide, Inches(0.3), y, Inches(0.05), row_h - Inches(0.1), color)
        # Label
        _add_text(slide, label,
                  Inches(0.45), y, Inches(9), Inches(0.22),
                  font_size=7, bold=True, color=SLATE)
        # Content
        _add_text(slide, text,
                  Inches(0.45), y + Inches(0.22), Inches(9.1), row_h - Inches(0.32),
                  font_size=9, color=DARK, wrap=True)
        y += row_h

    # Footer
    _add_text(slide, _vel_label(region.velocity),
              Inches(0.3), Inches(6.4), Inches(2.5), Inches(0.3),
              font_size=8, color=SLATE)
    _add_text(slide, region.threat_characterisation or "",
              Inches(3.0), Inches(6.4), Inches(4.0), Inches(0.3),
              font_size=8, color=SLATE)

    # Source quality badge
    sq = region.source_quality
    if sq and sq.get("total", 0) > 0:
        sq_text = f"Sources: {sq['tier_a']}A · {sq['tier_b']}B · {sq['tier_c']}C"
    elif region.top_sources:
        sq_text = "Sources: " + " · ".join(region.top_sources[:3])
    else:
        sq_text = ""
    if sq_text:
        _add_text(slide, sq_text,
                  Inches(0.3), Inches(6.8), Inches(9.4), Inches(0.28),
                  font_size=7, color=SLATE)

    # Speaker notes — full depth
    notes_parts = []
    if region.intel_bullets:
        notes_parts.append("INTEL FINDINGS:\n" + "\n".join(f"• {b}" for b in region.intel_bullets))
    if region.adversary_bullets:
        notes_parts.append("ADVERSARY ACTIVITY:\n" + "\n".join(f"• {b}" for b in region.adversary_bullets))
    if region.impact_bullets:
        notes_parts.append("IMPACT FOR AEROGRID:\n" + "\n".join(f"• {b}" for b in region.impact_bullets))
    if region.watch_bullets:
        notes_parts.append("WATCH FOR:\n" + "\n".join(f"• {b}" for b in region.watch_bullets))
    if region.action_bullets:
        notes_parts.append("RECOMMENDED ACTIONS:\n" + "\n".join(f"• {b}" for b in region.action_bullets))
    if notes_parts:
        slide.notes_slide.notes_text_frame.text = "\n\n".join(notes_parts)


def build_appendix(prs: Presentation, data: ReportData) -> None:
    slide = _add_blank_slide(prs)

    _add_rect(slide, 0, 0, W, Inches(0.65), BRAND)
    _add_text(slide, "Appendix", Inches(0.3), Inches(0.1),
              Inches(7), Inches(0.45), font_size=16, bold=True, color=WHITE)

    y = Inches(0.85)
    Y_MAX = Inches(7.1)  # slide height is 7.5in; stop adding rows near the bottom
    monitor = [r for r in data.regions if r.status == RegionStatus.MONITOR]
    clear   = [r for r in data.regions if r.status == RegionStatus.CLEAR]

    if monitor:
        _add_text(slide, "WATCH — MONITOR REGIONS", Inches(0.4), y,
                  Inches(9), Inches(0.3), font_size=9, bold=True, color=AMBER)
        y += Inches(0.35)
        for r in monitor:
            if y >= Y_MAX:
                break
            txt = f"{r.name}  ·  {r.scenario_match or 'No scenario'}  ·  Confidence: {r.confidence_label or '—'}  ·  {r.severity or '—'}"
            _add_text(slide, txt, Inches(0.6), y, Inches(9), Inches(0.35),
                      font_size=9, color=DARK)
            y += Inches(0.38)
        y += Inches(0.2)

    if clear and y < Y_MAX:
        _add_text(slide, "CLEAR REGIONS", Inches(0.4), y,
                  Inches(9), Inches(0.3), font_size=9, bold=True, color=GREEN)
        y += Inches(0.35)
        for r in clear:
            if y >= Y_MAX:
                break
            _add_text(slide, f"{r.name}  —  No active threat detected this cycle.",
                      Inches(0.6), y, Inches(9), Inches(0.35),
                      font_size=9, color=DARK)
            y += Inches(0.38)
        y += Inches(0.2)

    # Run metadata
    if y < Y_MAX:
        _add_text(slide, "RUN METADATA", Inches(0.4), y,
                  Inches(9), Inches(0.3), font_size=9, bold=True, color=SLATE)
        y += Inches(0.35)
    meta_rows = [
        ("Pipeline ID",    data.run_id),
        ("Timestamp",      data.timestamp),
        ("Escalated",      str(data.escalated_count)),
        ("Monitor",        str(data.monitor_count)),
        ("Clear",          str(data.clear_count)),
    ]
    for label, value in meta_rows:
        if y >= Y_MAX:
            break
        _add_text(slide, f"{label}:  {value}",
                  Inches(0.6), y, Inches(9), Inches(0.3),
                  font_size=9, color=DARK)
        y += Inches(0.3)


# ── Public API ─────────────────────────────────────────────────────────────────

def build_presentation(data: ReportData) -> Presentation:
    _ensure_base_pptx()
    prs = Presentation(str(BASE_PPTX))
    prs.slide_width  = W
    prs.slide_height = H

    build_cover(prs, data)
    build_exec_summary(prs, data)
    for region in (r for r in data.regions if r.status == RegionStatus.ESCALATED):
        build_region(prs, region)
    build_appendix(prs, data)

    return prs


def export(output_path: str = DEFAULT_OUT, output_dir: str = "output") -> None:
    data = build(output_dir=output_dir)
    prs  = build_presentation(data)
    prs.save(output_path)
    print(f"PPTX exported: {output_path}")


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_OUT
    export(out)
