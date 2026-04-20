"""
CRQ Deliverable -- PPTX builder (v2)
=====================================

Implements CRQ Design System v1.0 (tools/design/).
PDF (build_pdf.py) is the design reference -- this file ports it slide-by-slide to 16:9 PPTX.

v2 notes:
  - Scatter chart: native python-pptx shapes (no Matplotlib)
  - All data from shared brief_data.py
  - Token ramp, category circles, pill-combo from design system
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.path.insert(0, str(Path(__file__).parent))
from brief_data import BRIEF_DATA  # noqa: E402

# -- Design tokens (CRQ Design System v1.0) ------------------------------------

WHITE        = RGBColor(0xff, 0xff, 0xff)
PAPER        = RGBColor(0xfb, 0xfb, 0xfd)
INK_900      = RGBColor(0x0e, 0x15, 0x30)
INK_700      = RGBColor(0x1b, 0x2d, 0x6e)
INK_500      = RGBColor(0x4a, 0x52, 0x75)
INK_400      = RGBColor(0x6e, 0x75, 0x96)
INK_300      = RGBColor(0xa8, 0xad, 0xc4)
INK_200      = RGBColor(0xd6, 0xd9, 0xe5)
INK_150      = RGBColor(0xe4, 0xe6, 0xef)
INK_100      = RGBColor(0xee, 0xf0, 0xf6)
INK_50       = RGBColor(0xf5, 0xf6, 0xfa)

BRAND_NAVY   = INK_700
BRAND_BRIGHT = RGBColor(0x2e, 0x5b, 0xff)

SEV = {
    "CRITICAL": RGBColor(0xd1, 0x24, 0x2f),
    "HIGH":     RGBColor(0xcf, 0x6f, 0x12),
    "MEDIUM":   RGBColor(0x9a, 0x73, 0x00),
    "MONITOR":  RGBColor(0x09, 0x69, 0xda),
}

SEV_TINT = {
    "CRITICAL": RGBColor(0xfb, 0xe9, 0xeb),
    "HIGH":     RGBColor(0xfb, 0xee, 0xdf),
    "MEDIUM":   RGBColor(0xf6, 0xee, 0xd6),
    "MONITOR":  RGBColor(0xe1, 0xed, 0xfb),
}

CAT_COLOR = BRAND_NAVY
FONT = "IBM Plex Sans"


# ── Helper utilities ───────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=11, bold=False, color=None,
             family=FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.35, italic=False):
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = family
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    return box, tf, p, run


def add_rect(slide, x, y, w, h, fill, *, line_color=None, line_w=0.5):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_dot(slide, x, y, diameter, color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y), Inches(diameter), Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_pill(slide, x, y, w, h, fill_color, label, *, text_color=None, label_size=10):
    from pptx.enum.shapes import MSO_SHAPE
    if text_color is None:
        text_color = WHITE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(label_size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return shape


def add_category_circle(slide, x, y, diameter, letter, *, letter_size=11):
    """All-navy category circle per design spec P·03."""
    add_dot(slide, x, y, diameter, CAT_COLOR)
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(diameter), Inches(diameter)
    )
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = letter
    run.font.name = FONT
    run.font.size = Pt(letter_size)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_footer(slide, *, brief="GLOBAL CYCLE BRIEF · v2.0", page=1, total=12):
    y = 7.2
    add_text(slide, 0.5, y, 5, 0.22, brief, size=8, color=INK_400)
    add_text(slide, 0.5, y, 12.333, 0.22, "CRQ COMMAND CENTER",
             size=8, color=INK_400, align=PP_ALIGN.CENTER)
    add_text(slide, 7.333, y, 5.5, 0.22, f"{page:02d} / {total:02d}",
             size=8, color=INK_400, align=PP_ALIGN.RIGHT)


def add_brand_mark(slide, x, y):
    add_text(slide, x, y, 6, 0.28, "CRQ COMMAND CENTER",
             size=10, bold=True, color=BRAND_NAVY)
    add_text(slide, x, y + 0.26, 6, 0.22,
             "Cyber Risk Quantification · Threat Intelligence",
             size=8.5, color=INK_400)


# ── Slide layouts ──────────────────────────────────────────────────────────

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAPER)
    add_brand_mark(slide, 0.5, 0.5)

    add_text(slide, 0.5, 3.1, 12, 0.32,
             "GLOBAL THREAT BRIEF", size=11, bold=True, color=INK_400)
    add_text(slide, 0.5, 3.5, 12, 1.65,
             BRIEF_DATA["cover_title"], size=70, bold=True, color=BRAND_NAVY, line_spacing=1.05)
    add_text(slide, 0.5, 5.05, 12, 0.5,
             BRIEF_DATA["cover_subtitle"], size=22, color=INK_500, line_spacing=1.3)

    add_rect(slide, 0.5, 6.72, 12.333, 0.012, INK_200)
    meta_y = 6.87
    add_text(slide, 0.5,  meta_y, 4, 0.2, "PUBLISHED",   size=8, bold=True, color=INK_400)
    add_text(slide, 0.5,  meta_y + 0.2, 4, 0.28, BRIEF_DATA["published_date"], size=11, bold=True, color=INK_900)
    add_text(slide, 5.0,  meta_y, 4, 0.2, "PREPARED BY", size=8, bold=True, color=INK_400)
    add_text(slide, 5.0,  meta_y + 0.2, 4, 0.28, BRIEF_DATA["prepared_by"], size=11, bold=True, color=INK_900)
    add_text(slide, 10.0, meta_y, 3, 0.2, "VERSION",     size=8, bold=True, color=INK_400)
    add_text(slide, 10.0, meta_y + 0.2, 3, 0.28, BRIEF_DATA["version"].split(" ")[0], size=11, bold=True, color=INK_900)


def build_exec_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAPER)

    add_text(slide, 0.5, 0.5, 12, 0.24,
             f"EXECUTIVE SUMMARY  ·  CYCLE {BRIEF_DATA['cycle']}",
             size=9, bold=True, color=INK_400)
    add_text(slide, 0.5, 0.84, 12.333, 0.85,
             BRIEF_DATA["exec_top_line"],
             size=17, color=BRAND_NAVY, line_spacing=1.3)
    add_rect(slide, 0.5, 1.9, 12.333, 0.012, INK_700)

    add_text(slide, 0.5, 2.05, 7.5, 0.22, "KEY DEVELOPMENTS", size=9, bold=True, color=INK_400)

    ey = 2.36
    for e in BRIEF_DATA["key_developments"]:
        sev = e["severity"]
        is_new = e["is_new"]
        region = e["region"]
        headline = e["headline"]
        impact = e["impact"]

        card_h = 0.82 if impact else 0.52
        add_rect(slide, 0.5, ey, 7.5, card_h, WHITE,
                 line_color=INK_150, line_w=0.75)

        pill_label = f"{sev}  ·  {region}" + ("  ·  NEW" if is_new else "")
        add_pill(slide, 0.62, ey + 0.07, 2.4, 0.2,
                 SEV[sev], pill_label, label_size=7.5)

        add_text(slide, 0.62, ey + 0.3, 7.2, 0.3,
                 headline, size=10.5, bold=True, color=INK_900, line_spacing=1.2)
        if impact:
            add_text(slide, 0.62, ey + 0.56, 7.2, 0.28,
                     impact, size=9, color=INK_500, line_spacing=1.3)
        ey += card_h + 0.04

    rx, rw = 8.3, 4.53
    add_text(slide, rx, 2.05, rw, 0.22, "ALSO TRACKING", size=9, bold=True, color=INK_400)
    ay = 2.36
    for a in BRIEF_DATA["also_tracking"]:
        add_dot(slide, rx, ay + 0.07, 0.14, SEV[a["severity"]])
        box = slide.shapes.add_textbox(Inches(rx+0.22), Inches(ay), Inches(rw-0.25), Inches(0.38))
        tf = box.text_frame; tf.margin_left = tf.margin_right = 0; tf.margin_top = tf.margin_bottom = 0; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.25
        r1 = p.add_run(); r1.text = a["region"]; r1.font.name = FONT; r1.font.size = Pt(9.5); r1.font.bold = True; r1.font.color.rgb = INK_500
        r2 = p.add_run(); r2.text = f"  -  {a['text']}"; r2.font.name = FONT; r2.font.size = Pt(9.5); r2.font.color.rgb = INK_900
        ay += 0.42

    watch_y = 3.65
    add_rect(slide, rx, watch_y, rw, 1.3, INK_100)
    add_text(slide, rx+0.12, watch_y+0.08, rw-0.2, 0.22, "WATCH NEXT", size=8.5, bold=True, color=INK_700)
    add_text(slide, rx+0.12, watch_y+0.32, rw-0.2, 1.0,
             BRIEF_DATA["watch_next"],
             size=9.5, color=INK_700, line_spacing=1.5)

    add_rect(slide, rx, 5.15, rw, 0.32, WHITE, line_color=INK_700, line_w=0.75)
    add_text(slide, rx+0.12, 5.18, rw-0.2, 0.26,
             f"CONFIDENCE: {BRIEF_DATA['confidence']}  ·  ADMIRALTY {BRIEF_DATA['admiralty']}",
             size=9, bold=True, color=INK_700, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page=2, total=7)


def build_scenario_page(prs, scenario, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAPER)

    number = scenario["number"]
    region = scenario["region"]

    add_text(slide, 0.5, 0.5, 12, 0.24,
             f"{region}  ·  SCENARIO {number} OF 3",
             size=9, bold=True, color=INK_400)
    add_text(slide, 0.5, 0.82, 12.333, 0.72,
             scenario["headline"], size=20, bold=True, color=BRAND_NAVY, line_spacing=1.15)

    add_text(slide, 0.5, 1.65, 5, 0.22, "KEY DRIVERS & ASSUMPTIONS", size=9, bold=True, color=INK_400)
    dy = 1.94
    for d in scenario["drivers"]:
        add_category_circle(slide, 0.5, dy, 0.3, d["letter"], letter_size=12)
        add_text(slide, 0.93, dy - 0.02, 5.5, 0.3, d["label"], size=11, bold=True, color=INK_900)
        add_text(slide, 0.93, dy + 0.22, 5.5, 0.28, f"-> {d['implication']}", size=9.5, color=INK_500)
        dy += 0.55

    type_colors = {"BASELINE": BRAND_NAVY, "PLAUSIBLE": SEV["HIGH"], "WILDCARD": SEV["CRITICAL"]}
    pill_color = type_colors.get(scenario["type"], BRAND_NAVY)
    add_pill(slide, 0.5, 3.68, 12.333, 0.34, pill_color,
             f"{scenario['type']}  ·  {scenario['type_subtitle']}", label_size=10)

    add_text(slide, 0.5, 4.12, 12.333, 0.3, scenario["bold_conclusion"],
             size=12, bold=True, color=INK_900)
    add_text(slide, 0.5, 4.44, 12.333, 0.72, scenario["body"],
             size=10, color=INK_900, line_spacing=1.45)

    add_pill(slide, 0.5, 5.28, 12.333, 0.28, BRAND_NAVY, "LIKELIHOOD", label_size=9.5)
    add_text(slide, 0.5, 5.6, 12.333, 0.24, scenario["likelihood_verdict"],
             size=10, bold=True, color=INK_900, line_spacing=1.3)

    col_y = 5.92; col_w = 6.05
    add_pill(slide, 0.5, col_y, col_w, 0.26, SEV["CRITICAL"], "BUSINESS IMPLICATIONS", label_size=9.5)
    add_pill(slide, 6.78, col_y, col_w, 0.26, SEV["HIGH"], "PRIORITY ACTIONS  ·  0-90 DAYS", label_size=9.5)

    iy = col_y + 0.3
    for i in scenario["implications"]:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(iy), Inches(col_w-0.1), Inches(0.19))
        tf = box.text_frame; tf.margin_left = tf.margin_right = 0; tf.margin_top = tf.margin_bottom = 0; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▸  "; r0.font.name = FONT; r0.font.size = Pt(7); r0.font.color.rgb = INK_300
        r1 = p.add_run(); r1.text = f"{i['label']}: "; r1.font.name = FONT; r1.font.size = Pt(9); r1.font.bold = True; r1.font.color.rgb = INK_900
        r2 = p.add_run(); r2.text = i["text"]; r2.font.name = FONT; r2.font.size = Pt(9); r2.font.color.rgb = INK_500
        iy += 0.16

    ay = col_y + 0.3
    for a in scenario["actions"]:
        box = slide.shapes.add_textbox(Inches(6.83), Inches(ay), Inches(col_w-0.1), Inches(0.19))
        tf = box.text_frame; tf.margin_left = tf.margin_right = 0; tf.margin_top = tf.margin_bottom = 0; tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.15
        r0 = p.add_run(); r0.text = "▸  "; r0.font.name = FONT; r0.font.size = Pt(7); r0.font.color.rgb = INK_300
        r1 = p.add_run(); r1.text = a; r1.font.name = FONT; r1.font.size = Pt(9); r1.font.color.rgb = INK_900
        ay += 0.16

    add_text(slide, 0.5, 7.1, 12.333, 0.18, scenario["evidence"],
             size=8, color=INK_400, italic=True)
    add_footer(slide, page=page, total=7)


def build_matrix(prs):
    """Matrix slide: full-width python-pptx scatter chart matching PDF design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAPER)
    d = BRIEF_DATA["matrix"]

    # Header
    add_text(slide, 0.5, 0.50, 12, 0.24,
             "GLOBAL  ·  SCENARIO PORTFOLIO", size=9, bold=True, color=INK_400)
    add_text(slide, 0.5, 0.80, 12.333, 0.65,
             d["headline"], size=20, bold=True, color=BRAND_NAVY, line_spacing=1.15)
    add_rect(slide, 0.5, 1.58, 12.333, 0.012, INK_700)
    add_text(slide, 0.5, 1.70, 12, 0.22, "BOTTOM LINE", size=9, bold=True, color=INK_400)
    add_text(slide, 0.5, 1.95, 12.333, 0.55,
             d["bottom_line"], size=11, color=INK_900, line_spacing=1.45)

    # Scatter chart container (white background with ink border)
    CONT_L, CONT_T, CONT_W, CONT_H = 0.50, 2.60, 12.333, 3.85
    add_rect(slide, CONT_L, CONT_T, CONT_W, CONT_H, WHITE,
             line_color=INK_200, line_w=0.75)

    # Canvas coordinates (inset from container)
    PAD_L, PAD_T, PAD_R, PAD_B = 0.65, 0.45, 0.35, 0.50
    CL = CONT_L + PAD_L
    CT = CONT_T + PAD_T
    CR = CONT_L + CONT_W - PAD_R
    CB = CONT_T + CONT_H - PAD_B
    CW = CR - CL
    CH = CB - CT

    MX = CL + CW / 2
    MY = CT + CH / 2

    # Top-right quadrant tint
    add_rect(slide, MX, CT, CR - MX, MY - CT,
             SEV_TINT["CRITICAL"], line_color=None)

    # Axis crosshair lines
    add_rect(slide, MX - 0.005, CT, 0.01, CH, INK_700)
    add_rect(slide, CL, MY - 0.005, CW, 0.01, INK_700)

    # Axis labels
    add_text(slide, CL, CB + 0.07, 2.5, 0.22,
             "LOW LIKELIHOOD", size=8, bold=True, color=INK_400)
    add_text(slide, CR - 2.5, CB + 0.07, 2.5, 0.22,
             "HIGH LIKELIHOOD  ->", size=8, bold=True, color=INK_400,
             align=PP_ALIGN.RIGHT)

    # Y-axis rotated label
    ytb = slide.shapes.add_textbox(
        Inches(CONT_L + 0.05), Inches(CT), Inches(0.24), Inches(CH))
    ytb.rotation = 270
    ytf = ytb.text_frame
    ytf.margin_left = ytf.margin_right = ytf.margin_top = ytf.margin_bottom = 0
    ytf.vertical_anchor = MSO_ANCHOR.MIDDLE
    yp = ytf.paragraphs[0]
    yp.alignment = PP_ALIGN.CENTER
    yr = yp.add_run()
    yr.text = "LOW IMPACT  ^  HIGH IMPACT"
    yr.font.name = FONT
    yr.font.size = Pt(8)
    yr.font.bold = True
    yr.font.color.rgb = INK_400

    # Legend (top-right of container)
    LEG_ITEMS = [("CRITICAL", SEV["CRITICAL"]), ("HIGH", SEV["HIGH"]),
                 ("MEDIUM", SEV["MEDIUM"]), ("MONITOR", SEV["MONITOR"])]
    lx = CONT_L + CONT_W - 0.18
    for sev_name, sev_color in reversed(LEG_ITEMS):
        lw = len(sev_name) * 0.085 + 0.32
        lx -= lw
        add_dot(slide, lx, CONT_T + 0.12, 0.14, sev_color)
        add_text(slide, lx + 0.20, CONT_T + 0.10, lw - 0.22, 0.22,
                 sev_name, size=8, bold=True, color=INK_500)
        lx -= 0.12

    # Dots and labels
    DOT_D = 0.20
    for dot in d["dots"]:
        dx = CL + (dot["x"] / 100) * CW - DOT_D / 2
        dy = CB - (dot["y"] / 100) * CH - DOT_D / 2
        add_dot(slide, dx, dy, DOT_D, SEV[dot["severity"].upper()])
        add_text(slide, dx + DOT_D + 0.04, dy + 0.02, 2.2, 0.20,
                 dot["label"], size=7.5, bold=True, color=INK_700)

    # Footer row
    FY = 6.55
    add_text(slide, 0.5, FY, 1.4, 0.22,
             "READING THIS CHART", size=8, bold=True, color=INK_400)
    add_text(slide, 0.5, FY + 0.24, 7.0, 0.48,
             d["reading"], size=9.5, color=INK_900, line_spacing=1.45)

    add_text(slide, 7.8, FY, 2.0, 0.22,
             "SEVERITY", size=8, bold=True, color=INK_400)
    sy = FY + 0.24
    for sev_name, sev_color in LEG_ITEMS:
        add_dot(slide, 7.8, sy + 0.04, 0.15, sev_color)
        add_text(slide, 8.05, sy, 1.8, 0.22, sev_name, size=9.5, color=INK_900)
        sy += 0.26

    badge_label = f"ADMIRALTY {BRIEF_DATA['admiralty']}  ·  REFRESHED {BRIEF_DATA['refreshed']}"
    add_rect(slide, 9.9, 6.9, 2.933, 0.32, WHITE, line_color=INK_700, line_w=0.75)
    add_text(slide, 9.9, 6.9, 2.933, 0.32, badge_label,
             size=8, bold=True, color=INK_700,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page=6, total=7)


def build_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PAPER)

    add_text(slide, 0.5, 0.5, 12, 0.24, "METHODOLOGY  ·  SOURCES  ·  CHANGE LOG", size=9, bold=True, color=INK_400)
    add_text(slide, 0.5, 0.82, 12.333, 0.55,
             "How this brief was produced", size=20, bold=True, color=BRAND_NAVY)
    add_rect(slide, 0.5, 1.55, 12.333, 0.012, INK_700)

    c1x, cw = 0.5, 3.9
    add_pill(slide, c1x, 1.75, cw, 0.3, BRAND_NAVY, "PIPELINE METADATA", label_size=9.5)
    my = 2.2
    for label, val in BRIEF_DATA["methodology"]["meta"]:
        add_text(slide, c1x, my, cw, 0.2, label, size=8, bold=True, color=INK_400)
        add_text(slide, c1x, my + 0.2, cw, 0.28, val, size=10.5, color=INK_900, line_spacing=1.25)
        my += 0.58

    c2x = 4.75
    add_pill(slide, c2x, 1.75, cw, 0.3, BRAND_NAVY, "SOURCES (TOP, ABRIDGED)", label_size=9.5)
    sy = 2.2
    for grade, name in BRIEF_DATA["methodology"]["sources"]:
        add_pill(slide, c2x, sy + 0.04, 0.4, 0.22, BRAND_NAVY, grade, label_size=7.5)
        add_text(slide, c2x + 0.52, sy + 0.04, cw - 0.6, 0.28, name, size=10, color=INK_900)
        sy += 0.34

    c3x = 9.0
    add_pill(slide, c3x, 1.75, cw, 0.3, BRAND_NAVY, f"CHANGE vs. {BRIEF_DATA['previous_cycle']}", label_size=9.5)
    cy = 2.2
    for marker, text in BRIEF_DATA["methodology"]["changes"]:
        add_text(slide, c3x, cy, 0.28, 0.3, marker, size=11, bold=True, color=BRAND_NAVY)
        add_text(slide, c3x + 0.32, cy + 0.02, cw - 0.4, 0.38, text, size=10, color=INK_900, line_spacing=1.35)
        cy += 0.34

    add_rect(slide, 0.5, 6.42, 12.333, 0.012, INK_200)
    add_text(slide, 0.5, 6.55, 12.333, 0.5,
             "Full source list, IOC tables, and pipeline trace available as separate appendix artifacts. "
             "This brief is published per cycle (weekly) and is intended for internal stakeholders only.",
             size=9, color=INK_400, italic=True, line_spacing=1.5)

    add_footer(slide, page=7, total=7)


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    out_dir = Path(__file__).parent.parent / "output" / "deliverables"
    out_dir.mkdir(parents=True, exist_ok=True)

    print("Building deck...")
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_cover(prs)
    build_exec_summary(prs)
    for idx, scenario in enumerate(BRIEF_DATA["scenarios"], start=3):
        build_scenario_page(prs, scenario, page=idx)
    build_matrix(prs)
    build_methodology(prs)

    deck_path = out_dir / "CRQ-cycle-brief-v2.pptx"
    prs.save(deck_path)
    print(f"Saved: {deck_path}")


if __name__ == "__main__":
    main()
