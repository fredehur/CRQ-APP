"""TDD tests for the board PPTX redesign (export_pptx.py)."""
import pytest
from pptx import Presentation


# ── Helpers ────────────────────────────────────────────────────────────────────

def _all_text(pptx_path: str) -> str:
    """Extract all visible text from every slide."""
    prs = Presentation(pptx_path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            parts.append(run.text.strip())
    return " ".join(parts)


def _slide_text(pptx_path: str, slide_idx: int) -> str:
    """Extract visible text from a single slide by index."""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        parts.append(run.text.strip())
    return " ".join(parts)


def _fill_colors(pptx_path: str) -> set[str]:
    """Return hex strings of every solid fill color used in slide shapes."""
    prs = Presentation(pptx_path)
    colors: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                rgb = shape.fill.fore_color.rgb
                colors.add(str(rgb))
            except Exception:
                pass
    return colors


def _slide_count(pptx_path: str) -> int:
    return len(Presentation(pptx_path).slides)


# ── Tests ──────────────────────────────────────────────────────────────────────

def test_export_runs_without_error(mock_output, tmp_path):
    """export() completes without error and produces a non-trivial file."""
    out = tmp_path / "board_test.pptx"
    from tools.export_pptx import export
    export(output_path=str(out), output_dir=str(mock_output))
    assert out.exists(), "output file not created"
    assert out.stat().st_size > 5000, "file suspiciously small"


def test_cover_shows_plain_language_labels(mock_output, tmp_path):
    """Cover slide uses plain labels (Active Threats / Under Watch / No Active Threat),
    not pipeline enum values (ESCALATED / MONITOR / CLEAR)."""
    out = tmp_path / "board_test.pptx"
    from tools.export_pptx import export
    export(output_path=str(out), output_dir=str(mock_output))
    cover = _slide_text(str(out), 0)
    assert "Active Threats" in cover, f"'Active Threats' not found in cover: {cover}"
    assert "Under Watch" in cover, f"'Under Watch' not found in cover: {cover}"
    assert "No Active Threat" in cover, f"'No Active Threat' not found in cover: {cover}"
    # Enum values must NOT appear on slide face
    assert "ESCALATED" not in cover, "raw 'ESCALATED' found on cover"
    assert "MONITOR" not in cover, "raw 'MONITOR' found on cover"
    assert "CLEAR" not in cover, "raw 'CLEAR' found on cover"


def test_overview_slide_present(mock_output, tmp_path):
    """Slide 2 is the Intelligence Overview — contains section label."""
    out = tmp_path / "board_test.pptx"
    from tools.export_pptx import export
    export(output_path=str(out), output_dir=str(mock_output))
    overview = _slide_text(str(out), 1)
    assert "Intelligence Overview" in overview, (
        f"'Intelligence Overview' not in slide 2: {overview}"
    )


def test_deck_slide_count(mock_output, tmp_path):
    """With 2 escalated (different scenarios) + 1 monitor, deck has 5 slides:
    Cover + Overview + 2 threat slides + 1 watch list."""
    out = tmp_path / "board_test.pptx"
    from tools.export_pptx import export
    export(output_path=str(out), output_dir=str(mock_output))
    assert _slide_count(str(out)) == 5, (
        f"Expected 5 slides, got {_slide_count(str(out))}"
    )


def test_watch_list_present_for_monitor_regions(mock_output, tmp_path):
    """Watch list slide exists and contains 'Under Watch' + MED (the MONITOR region)."""
    out = tmp_path / "board_test.pptx"
    from tools.export_pptx import export
    export(output_path=str(out), output_dir=str(mock_output))
    full_text = _all_text(str(out))
    assert "Under Watch" in full_text, "'Under Watch' not found in deck"
    # MED is the MONITOR region in mock data
    last_slide = _slide_text(str(out), -1)
    assert "MED" in last_slide, f"MED not found in watch list slide: {last_slide}"


def test_no_red_amber_green_brand_colors(mock_output, tmp_path):
    """No RED, AMBER, GREEN, or BRAND NAVY fill colours appear in any slide shape."""
    out = tmp_path / "board_test.pptx"
    from tools.export_pptx import export
    export(output_path=str(out), output_dir=str(mock_output))
    forbidden = {"DC2626", "D97706", "16A34A", "104277"}
    found = _fill_colors(str(out)) & forbidden
    assert not found, f"Forbidden fill colours found in deck: {found}"
